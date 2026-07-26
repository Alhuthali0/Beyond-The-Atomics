import sys
from pptx import Presentation
from pdfminer.high_level import extract_text

def read_pptx(file_path):
    prs = Presentation(file_path)
    text = []
    for i, slide in enumerate(prs.slides):
        text.append(f"--- Slide {i+1} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def read_pdf(file_path):
    return extract_text(file_path)

if __name__ == "__main__":
    print("### PDF CONTENT ###")
    try:
        print(read_pdf("GP2 Presentation general guide notes.pdf"))
    except Exception as e:
        print(f"Error reading PDF: {e}")
    
    print("\n\n### PPTX CONTENT ###")
    try:
        print(read_pptx("Graduation project presentation.pptx"))
    except Exception as e:
        print(f"Error reading PPTX: {e}")
